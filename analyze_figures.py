# -*- coding: utf-8 -*-
"""
Analyze presentations for figure/stats presentation patterns
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os
import sys

# Fix encoding
sys.stdout.reconfigure(encoding='utf-8')

def emu_to_inches(emu):
    return emu / 914400

def analyze_presentation(pptx_path):
    """Analyze presentation for stats/figures patterns"""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error: {e}")
        return

    print(f"\n{'='*60}")
    print(f"FILE: {os.path.basename(pptx_path)}")
    print(f"{'='*60}")

    for slide_num, slide in enumerate(prs.slides, 1):
        texts_with_numbers = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Look for slides with numbers/percentages
                if text and any(char.isdigit() for char in text):
                    if any(x in text for x in ['%', '+', 'NOK', 'kr', 'M', 'B', 'mill', 'billion']):
                        texts_with_numbers.append(text[:150])
                    elif len(text) < 10 and any(char.isdigit() for char in text):
                        # Short text with numbers (likely a stat)
                        texts_with_numbers.append(text)

        if texts_with_numbers:
            print(f"\n--- SLIDE {slide_num} (Stats) ---")
            for t in texts_with_numbers[:8]:
                print(f"  [{t}]")

# Check key decks
decks = [
    r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Hafslund - Startuplab partnership Aug 25.pptx",
    r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Startuplab Founders Fund V presentation 28112023_LONG.pptx",
]

for deck in decks:
    if os.path.exists(deck):
        analyze_presentation(deck)
