"""
StartupLab PowerPoint Template Generator
Creates a branded presentation template with 10 master layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Brand Colors
COLORS = {
    'red': RGBColor(0xFF, 0x33, 0x33),
    'black': RGBColor(0x13, 0x14, 0x15),
    'deep_black': RGBColor(0x05, 0x05, 0x05),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'off_white': RGBColor(0xF7, 0xF7, 0xF7),
    'light_gray': RGBColor(0xF0, 0xF0, 0xF0),
    'muted': RGBColor(0x9B, 0xA1, 0xA5),
    'blue': RGBColor(0x1C, 0x64, 0xFF),
}

# Asset paths
ASSETS_DIR = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\brand-profile"
LOGO_WHITE = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_white.png")
LOGO_BLACK = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_black.png")
LOGO_RED = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_red.png")
SYMBOL_WHITE = os.path.join(ASSETS_DIR, "symbol", "png", "SL_symbol_white.png")
SYMBOL_BLACK = os.path.join(ASSETS_DIR, "symbol", "png", "SL_symbol_black.png")
SYMBOL_RED = os.path.join(ASSETS_DIR, "symbol", "png", "SL_symbol_red.png")

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_name="Inter", bold=False, color=COLORS['black'],
                 alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
    """Add a text box with specified formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

    tf.anchor = vertical
    return txBox

def add_placeholder_shape(slide, left, top, width, height, text="", color=COLORS['light_gray']):
    """Add a placeholder rectangle"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_shape_fill(shape, color)
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['muted']
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

    return shape

def create_template():
    """Create the full PowerPoint template"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]

    # ========================================
    # SLIDE 1: Title Slide (Red Background)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Red background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['red'])
    background.line.fill.background()

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(11), Inches(2),
        "PRESENTATION TITLE", font_size=60, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Subtitle
    add_text_box(
        slide, Inches(0.8), Inches(4.5), Inches(8), Inches(1),
        "Subtitle or date goes here", font_size=24, font_name="Inter",
        color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Logo (white on red)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(0.8), Inches(6.3), height=Inches(0.5))

    # ========================================
    # SLIDE 2: Title Slide (Dark Background)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['deep_black'])
    background.line.fill.background()

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(11), Inches(2),
        "PRESENTATION TITLE", font_size=60, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Subtitle
    add_text_box(
        slide, Inches(0.8), Inches(4.5), Inches(8), Inches(1),
        "Subtitle or date goes here", font_size=24, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(0.8), Inches(6.3), height=Inches(0.5))

    # ========================================
    # SLIDE 3: Section Divider
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Red background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['red'])
    background.line.fill.background()

    # Section number
    add_text_box(
        slide, Inches(0.8), Inches(2), Inches(2), Inches(1),
        "01", font_size=72, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Section title
    add_text_box(
        slide, Inches(0.8), Inches(3.2), Inches(10), Inches(1.5),
        "SECTION TITLE", font_size=48, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Symbol
    if os.path.exists(SYMBOL_WHITE):
        slide.shapes.add_picture(SYMBOL_WHITE, Inches(11.5), Inches(6), height=Inches(0.8))

    # ========================================
    # SLIDE 4: Content Slide (White)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # White background (default)

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
        "Slide Title", font_size=36, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Body text
    add_text_box(
        slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
        "Body text goes here. Use Inter font at 17-19px for optimal readability.\n\n• Bullet point one\n• Bullet point two\n• Bullet point three",
        font_size=18, font_name="Inter", color=COLORS['black'],
        alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 5: Image Left + Text Right
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Image placeholder (left half)
    add_placeholder_shape(
        slide, Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT,
        "IMAGE", COLORS['light_gray']
    )

    # Title (right)
    add_text_box(
        slide, Inches(7), Inches(1.5), Inches(5.5), Inches(1),
        "Slide Title", font_size=32, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Body text (right)
    add_text_box(
        slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5),
        "Supporting text goes here. This layout works well for case studies, features, or storytelling.",
        font_size=18, font_name="Inter", color=COLORS['black'],
        alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(7), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 6: Text Left + Image Right
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Title (left)
    add_text_box(
        slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1),
        "Slide Title", font_size=32, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Body text (left)
    add_text_box(
        slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(3.5),
        "Supporting text goes here. This layout works well for case studies, features, or storytelling.",
        font_size=18, font_name="Inter", color=COLORS['black'],
        alignment=PP_ALIGN.LEFT
    )

    # Image placeholder (right half)
    add_placeholder_shape(
        slide, Inches(6.833), Inches(0), Inches(6.5), SLIDE_HEIGHT,
        "IMAGE", COLORS['light_gray']
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 7: Data/Stats Slide
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
        "Key Metrics", font_size=36, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Stat 1
    add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(1.2),
        "250+", font_size=72, font_name="Arial Black",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )
    add_text_box(
        slide, Inches(0.8), Inches(3.6), Inches(3.5), Inches(0.8),
        "Startups supported", font_size=18, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Stat 2
    add_text_box(
        slide, Inches(5), Inches(2.2), Inches(3.5), Inches(1.2),
        "15B+", font_size=72, font_name="Arial Black",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )
    add_text_box(
        slide, Inches(5), Inches(3.6), Inches(3.5), Inches(0.8),
        "NOK raised by portfolio", font_size=18, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Stat 3
    add_text_box(
        slide, Inches(9.2), Inches(2.2), Inches(3.5), Inches(1.2),
        "10+", font_size=72, font_name="Arial Black",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )
    add_text_box(
        slide, Inches(9.2), Inches(3.6), Inches(3.5), Inches(0.8),
        "Industry programs", font_size=18, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 8: Quote Slide
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Off-white background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['off_white'])
    background.line.fill.background()

    # Red quote mark
    add_text_box(
        slide, Inches(0.8), Inches(1.5), Inches(1), Inches(1),
        "\u201C", font_size=120, font_name="Georgia",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )

    # Quote text
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(10), Inches(2.5),
        "Quote text goes here. Make it impactful and memorable.",
        font_size=32, font_name="Georgia",
        color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Attribution
    add_text_box(
        slide, Inches(0.8), Inches(5.3), Inches(10), Inches(0.5),
        "— Name, Title", font_size=18, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 9: Team/Profile Slide
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
        "Our Team", font_size=36, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Team member placeholders (4 columns)
    for i in range(4):
        x = Inches(0.8 + i * 3.1)

        # Photo placeholder
        add_placeholder_shape(
            slide, x, Inches(1.8), Inches(2.5), Inches(2.5),
            "PHOTO", COLORS['light_gray']
        )

        # Name
        add_text_box(
            slide, x, Inches(4.5), Inches(2.5), Inches(0.5),
            "Name Here", font_size=18, font_name="Arial",
            bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
        )

        # Title
        add_text_box(
            slide, x, Inches(5), Inches(2.5), Inches(0.5),
            "Job Title", font_size=14, font_name="Inter",
            color=COLORS['muted'], alignment=PP_ALIGN.LEFT
        )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 10: Two Column Content
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Title
    add_text_box(
        slide, Inches(0.8), Inches(0.6), Inches(11), Inches(1),
        "Two Column Title", font_size=36, font_name="Arial",
        bold=True, color=COLORS['black'], alignment=PP_ALIGN.LEFT
    )

    # Left column title
    add_text_box(
        slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.6),
        "Column One", font_size=24, font_name="Arial",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )

    # Left column content
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5),
        "Content for the first column goes here.\n\n• Point one\n• Point two\n• Point three",
        font_size=16, font_name="Inter", color=COLORS['black'],
        alignment=PP_ALIGN.LEFT
    )

    # Right column title
    add_text_box(
        slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.6),
        "Column Two", font_size=24, font_name="Arial",
        bold=True, color=COLORS['red'], alignment=PP_ALIGN.LEFT
    )

    # Right column content
    add_text_box(
        slide, Inches(7), Inches(2.5), Inches(5.5), Inches(3.5),
        "Content for the second column goes here.\n\n• Point one\n• Point two\n• Point three",
        font_size=16, font_name="Inter", color=COLORS['black'],
        alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_BLACK):
        slide.shapes.add_picture(LOGO_BLACK, Inches(0.8), Inches(6.7), height=Inches(0.4))

    # ========================================
    # SLIDE 11: Closing/CTA Slide
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['deep_black'])
    background.line.fill.background()

    # CTA headline
    add_text_box(
        slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.5),
        "READY TO GO\nFURTHER FASTER?", font_size=48, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Contact info
    add_text_box(
        slide, Inches(0.8), Inches(4.5), Inches(6), Inches(1),
        "contact@startuplab.no  |  startuplab.no", font_size=18, font_name="Inter",
        color=COLORS['muted'], alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(0.8), Inches(6.3), height=Inches(0.5))

    # Symbol (large, right side)
    if os.path.exists(SYMBOL_WHITE):
        slide.shapes.add_picture(SYMBOL_WHITE, Inches(10), Inches(2), height=Inches(3))

    # ========================================
    # SLIDE 12: Thank You Slide
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Red background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS['red'])
    background.line.fill.background()

    # Thank you text
    add_text_box(
        slide, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
        "THANK YOU", font_size=72, font_name="Arial Black",
        bold=True, color=COLORS['white'], alignment=PP_ALIGN.LEFT
    )

    # Logo
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(0.8), Inches(6.3), height=Inches(0.5))

    # Save the presentation
    output_path = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\StartupLab_Template.pptx"
    prs.save(output_path)
    print(f"Template saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_template()
