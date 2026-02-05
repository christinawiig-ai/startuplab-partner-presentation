"""
Extend the existing StartupLab template with additional text-heavy layouts
Uses the official template as base and adds new slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os
import shutil

# Brand Colors
COLORS = {
    'red': RGBColor(0xFF, 0x33, 0x33),
    'black': RGBColor(0x13, 0x14, 0x15),
    'deep_black': RGBColor(0x05, 0x05, 0x05),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'off_white': RGBColor(0xF7, 0xF7, 0xF7),
    'light_gray': RGBColor(0xE5, 0xE5, 0xE5),
    'muted': RGBColor(0x9B, 0xA1, 0xA5),
}

# Fonts - Using Replica as primary (with fallback to Rubik which is in template)
FONT_HEADLINE = "Replica LL"  # Or "Rubik Medium" as fallback
FONT_BODY = "Replica LL"      # Or "Rubik" as fallback
FONT_HEADLINE_FALLBACK = "Rubik Medium"
FONT_BODY_FALLBACK = "Rubik"

# Template dimensions (from analysis)
SLIDE_WIDTH = Inches(10.0)
SLIDE_HEIGHT = Inches(5.62)

# Standard positions (from analysis)
MARGIN_LEFT = Inches(0.3)
MARGIN_RIGHT = Inches(0.3)
LOGO_POS = (Inches(0.1), Inches(5.2))
LOGO_SIZE = Inches(0.3)

# Asset paths
ASSETS_DIR = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\brand-profile"
LOGO_WHITE = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_white.png")
LOGO_BLACK = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_black.png")
SYMBOL_WHITE = os.path.join(ASSETS_DIR, "symbol", "png", "SL_symbol_white.png")
SYMBOL_BLACK = os.path.join(ASSETS_DIR, "symbol", "png", "SL_symbol_black.png")


def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=17,
                 font_name=FONT_BODY, bold=False, color=COLORS['black'],
                 alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP,
                 all_caps=False, line_spacing=1.2):
    """Add a text box with specified formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text.upper() if all_caps else text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment

    # Line spacing
    p.line_spacing = line_spacing

    tf.anchor = vertical
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=15,
                    font_name=FONT_BODY, color=COLORS['black'], bullet_color=COLORS['red']):
    """Add a text box with bullet points"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.level = 0
        p.line_spacing = 1.5

        # Add bullet
        p.bullet = True

    return txBox


def add_logo(slide, light_bg=True):
    """Add logo to standard position"""
    logo_path = LOGO_BLACK if light_bg else LOGO_WHITE
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, LOGO_POS[0], LOGO_POS[1], height=LOGO_SIZE)


def create_extended_template():
    """Create extended template starting from existing one"""

    # Source template
    source_path = r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Presentation Template (Updated_ 24.04.24).pptx"
    output_path = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\StartupLab_Extended_Template.pptx"

    # Copy the original template
    shutil.copy(source_path, output_path)

    # Open and extend
    prs = Presentation(output_path)
    blank_layout = prs.slide_layouts[10]  # BLANK layout

    print(f"Starting with {len(prs.slides)} slides from original template")

    # ========================================
    # NEW SLIDE: Full Text Content (Light)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Headline
    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "FULL WIDTH TEXT SLIDE", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    # Subheadline
    add_text_box(
        slide, MARGIN_LEFT, Inches(1.1), Inches(9.4), Inches(0.4),
        "Optional subheadline for additional context", font_size=15,
        font_name=FONT_BODY_FALLBACK, color=COLORS['muted']
    )

    # Body text - full width
    add_text_box(
        slide, MARGIN_LEFT, Inches(1.8), Inches(9.4), Inches(3.0),
        "Use this layout when you need to present longer text content without images. "
        "Keep paragraphs concise and use line breaks between sections for readability.\n\n"
        "The full width allows for detailed explanations, process descriptions, or "
        "comprehensive information that needs more space than a split layout provides.",
        font_size=17, font_name=FONT_BODY_FALLBACK, color=COLORS['black'],
        line_spacing=1.4
    )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Bullet Points (Light)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "BULLET POINT SLIDE", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    add_bullet_text(
        slide, MARGIN_LEFT, Inches(1.4), Inches(9.4), Inches(3.5),
        [
            "First key point with supporting detail that explains the concept clearly",
            "Second point that builds on the first and adds new information",
            "Third point with additional context or examples",
            "Fourth point to round out the main message",
            "Optional fifth point if needed for completeness"
        ],
        font_size=17, font_name=FONT_BODY_FALLBACK
    )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Two Column Text (Light)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "TWO COLUMN COMPARISON", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    # Left column header
    add_text_box(
        slide, MARGIN_LEFT, Inches(1.3), Inches(4.5), Inches(0.4),
        "COLUMN ONE", font_size=15, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['red'], all_caps=True
    )

    # Left column body
    add_text_box(
        slide, MARGIN_LEFT, Inches(1.8), Inches(4.5), Inches(3.0),
        "Content for the left column goes here. Use this layout for comparisons, "
        "before/after scenarios, or presenting two related topics side by side.\n\n"
        "Keep the text balanced between columns for visual harmony.",
        font_size=15, font_name=FONT_BODY_FALLBACK, color=COLORS['black'],
        line_spacing=1.3
    )

    # Right column header
    add_text_box(
        slide, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.4),
        "COLUMN TWO", font_size=15, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['red'], all_caps=True
    )

    # Right column body
    add_text_box(
        slide, Inches(5.2), Inches(1.8), Inches(4.5), Inches(3.0),
        "Content for the right column. This side should complement the left column "
        "and together tell a complete story.\n\n"
        "Consider using bullet points in one or both columns if listing items.",
        font_size=15, font_name=FONT_BODY_FALLBACK, color=COLORS['black'],
        line_spacing=1.3
    )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Three Column Layout
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "THREE KEY POINTS", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    col_width = Inches(2.9)
    for i, (num, title, desc) in enumerate([
        ("01", "FIRST POINT", "Description of the first key concept or feature that you want to highlight."),
        ("02", "SECOND POINT", "Description of the second important element that builds your narrative."),
        ("03", "THIRD POINT", "Description of the third element that completes your message.")
    ]):
        x = MARGIN_LEFT + (i * Inches(3.2))

        # Number
        add_text_box(
            slide, x, Inches(1.4), Inches(1), Inches(0.6),
            num, font_size=36, font_name=FONT_HEADLINE_FALLBACK,
            bold=True, color=COLORS['red']
        )

        # Title
        add_text_box(
            slide, x, Inches(2.1), col_width, Inches(0.4),
            title, font_size=14, font_name=FONT_HEADLINE_FALLBACK,
            bold=True, color=COLORS['black'], all_caps=True
        )

        # Description
        add_text_box(
            slide, x, Inches(2.6), col_width, Inches(2.0),
            desc, font_size=13, font_name=FONT_BODY_FALLBACK,
            color=COLORS['black'], line_spacing=1.3
        )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Process/Timeline
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "PROCESS OVERVIEW", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    # Process steps (4 columns)
    steps = [
        ("1", "DISCOVER", "Understand the problem and gather requirements"),
        ("2", "DESIGN", "Create solutions and validate with stakeholders"),
        ("3", "DEVELOP", "Build and test the solution iteratively"),
        ("4", "DELIVER", "Launch, measure, and optimize continuously")
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = MARGIN_LEFT + (i * Inches(2.4))

        # Circle with number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.6), Inches(0.6), Inches(0.6))
        set_shape_fill(circle, COLORS['red'])
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Arrow (except for last)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.5), Inches(1.75), Inches(0.7), Inches(0.3))
            set_shape_fill(arrow, COLORS['light_gray'])
            arrow.line.fill.background()

        # Title
        add_text_box(
            slide, x, Inches(2.5), Inches(2.2), Inches(0.4),
            title, font_size=13, font_name=FONT_HEADLINE_FALLBACK,
            bold=True, color=COLORS['black'], all_caps=True, alignment=PP_ALIGN.CENTER
        )

        # Description
        add_text_box(
            slide, x, Inches(3.0), Inches(2.2), Inches(1.5),
            desc, font_size=12, font_name=FONT_BODY_FALLBACK,
            color=COLORS['muted'], alignment=PP_ALIGN.CENTER, line_spacing=1.2
        )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Quote (Dark Background)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(bg, COLORS['deep_black'])
    bg.line.fill.background()

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_LEFT, Inches(1.5), Inches(0.08), Inches(2.5))
    set_shape_fill(line, COLORS['red'])
    line.line.fill.background()

    # Quote
    add_text_box(
        slide, Inches(0.6), Inches(1.5), Inches(8.5), Inches(2.0),
        '"This is a powerful quote that captures the essence of your message. '
        'Keep it concise and impactful."',
        font_size=28, font_name=FONT_BODY_FALLBACK,
        color=COLORS['white'], line_spacing=1.3
    )

    # Attribution
    add_text_box(
        slide, Inches(0.6), Inches(3.8), Inches(8.5), Inches(0.5),
        "— Name, Title at Company",
        font_size=15, font_name=FONT_BODY_FALLBACK,
        color=COLORS['muted']
    )

    add_logo(slide, light_bg=False)

    # ========================================
    # NEW SLIDE: Key Takeaways
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(9.4), Inches(0.6),
        "KEY TAKEAWAYS", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    takeaways = [
        "First key takeaway that summarizes an important point from the presentation",
        "Second takeaway that highlights another crucial insight or action item",
        "Third takeaway that reinforces the main message and call to action"
    ]

    for i, text in enumerate(takeaways):
        y = Inches(1.4 + i * 1.2)

        # Checkmark or number box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_LEFT, y, Inches(0.5), Inches(0.5))
        set_shape_fill(box, COLORS['red'])
        box.line.fill.background()
        box.adjustments[0] = 0.2  # Corner radius
        tf = box.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Takeaway text
        add_text_box(
            slide, Inches(1.0), y, Inches(8.5), Inches(0.8),
            text, font_size=17, font_name=FONT_BODY_FALLBACK,
            color=COLORS['black'], vertical=MSO_ANCHOR.MIDDLE, line_spacing=1.2
        )

    add_logo(slide, light_bg=True)

    # ========================================
    # NEW SLIDE: Contact/CTA (Dark)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_shape_fill(bg, COLORS['deep_black'])
    bg.line.fill.background()

    # Headline
    add_text_box(
        slide, MARGIN_LEFT, Inches(1.5), Inches(6), Inches(1.2),
        "LET'S BUILD\nSOMETHING GREAT", font_size=40, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['white'], all_caps=True
    )

    # Contact details
    add_text_box(
        slide, MARGIN_LEFT, Inches(3.2), Inches(6), Inches(0.4),
        "name@startuplab.no", font_size=17, font_name=FONT_BODY_FALLBACK,
        color=COLORS['muted']
    )

    add_text_box(
        slide, MARGIN_LEFT, Inches(3.6), Inches(6), Inches(0.4),
        "startuplab.no", font_size=17, font_name=FONT_BODY_FALLBACK,
        color=COLORS['red']
    )

    # Symbol on right
    if os.path.exists(SYMBOL_WHITE):
        slide.shapes.add_picture(SYMBOL_WHITE, Inches(7.5), Inches(1.5), height=Inches(2.5))

    add_logo(slide, light_bg=False)

    # ========================================
    # NEW SLIDE: Text + Small Image (Right)
    # ========================================
    slide = prs.slides.add_slide(blank_layout)

    add_text_box(
        slide, MARGIN_LEFT, Inches(0.4), Inches(5.5), Inches(0.6),
        "TEXT-HEAVY WITH IMAGE", font_size=28, font_name=FONT_HEADLINE_FALLBACK,
        bold=True, color=COLORS['black'], all_caps=True
    )

    add_text_box(
        slide, MARGIN_LEFT, Inches(1.2), Inches(5.5), Inches(3.5),
        "This layout accommodates more text while still including a visual element. "
        "Use when you need to explain a concept in detail but want to maintain visual interest.\n\n"
        "Key points to cover:\n"
        "• First important detail\n"
        "• Second important detail\n"
        "• Third important detail\n\n"
        "Additional context or explanation can go here to fully develop the idea.",
        font_size=15, font_name=FONT_BODY_FALLBACK, color=COLORS['black'],
        line_spacing=1.3
    )

    # Image placeholder (smaller, right side)
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(0.8), Inches(3.5), Inches(3.5)
    )
    set_shape_fill(img_placeholder, COLORS['light_gray'])
    img_placeholder.line.fill.background()
    tf = img_placeholder.text_frame
    tf.paragraphs[0].text = "IMAGE"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLORS['muted']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

    add_logo(slide, light_bg=True)

    # Save
    prs.save(output_path)
    print(f"\nExtended template saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    create_extended_template()
