"""
Analyze corporate collaboration presentations for patterns
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def emu_to_inches(emu):
    return emu / 914400

def analyze_presentation(pptx_path):
    """Analyze presentation for collab/figures patterns"""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening {pptx_path}: {e}")
        return

    print(f"\n{'='*60}")
    print(f"ANALYZING: {os.path.basename(pptx_path)}")
    print(f"{'='*60}")
    print(f"Slides: {len(prs.slides)}")

    for slide_num, slide in enumerate(prs.slides, 1):
        # Look for interesting patterns
        texts = []
        has_numbers = False
        has_percentages = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text[:100])
                    # Check for numbers/stats
                    if any(char.isdigit() for char in text):
                        if '%' in text or 'NOK' in text or 'kr' in text or '+' in text:
                            has_numbers = True
                        if any(x in text.lower() for x in ['startup', 'partner', 'collab', 'corp', 'invest']):
                            has_numbers = True

        # Only print slides with interesting content
        combined = ' '.join(texts).lower()
        if any(x in combined for x in ['partner', 'startup', 'corporate', 'collab', 'pilot', 'poc',
                                        'invest', 'value', 'benefit', '%', 'nok', 'million', 'billion']):
            print(f"\n--- SLIDE {slide_num} ---")
            for t in texts[:5]:
                if t:
                    print(f"  {t[:80]}{'...' if len(t) > 80 else ''}")

# Analyze relevant presentations
decks = [
    r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Startuplab Collaboration Framework v3.0.pptx",
    r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Hafslund - Startuplab partnership Aug 25.pptx",
    r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Corp Partner Check In Aug_Sept 25.pptx",
    r"C:\Users\Christina\Code\Work\Prospecting tool\knowledge-base\cases\Corp Startup Collaboration Examples.pptx",
    r"C:\Users\Christina\Code\Work\notion_export\Startuplabs_Norconsult_281124_(1).pptx",
]

for deck in decks:
    if os.path.exists(deck):
        analyze_presentation(deck)
    else:
        print(f"Not found: {deck}")
