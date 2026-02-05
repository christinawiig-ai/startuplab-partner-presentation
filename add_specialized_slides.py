# -*- coding: utf-8 -*-
"""
Add specialized slides for corporate collaboration and key figures
Based on analysis of existing StartupLab presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand Colors
COLORS = {
    'red': RGBColor(0xFF, 0x33, 0x33),
    'black': RGBColor(0x13, 0x14, 0x15),
    'deep_black': RGBColor(0x05, 0x05, 0x05),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'off_white': RGBColor(0xF7, 0xF7, 0xF7),
    'light_gray': RGBColor(0xE8, 0xE8, 0xE8),
    'muted': RGBColor(0x9B, 0xA1, 0xA5),
}

FONT_HEADLINE = "Rubik Medium"
FONT_BODY = "Rubik"
SLIDE_WIDTH = Inches(10.0)
SLIDE_HEIGHT = Inches(5.62)
MARGIN = Inches(0.3)

ASSETS_DIR = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\brand-profile"
LOGO_WHITE = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_white.png")
LOGO_BLACK = os.path.join(ASSETS_DIR, "logo", "png", "SL_signature_black.png")


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=17, font=FONT_BODY,
             bold=False, color=COLORS['black'], align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, caps=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text.upper() if caps else text
    p.font.size = Pt(size)
    p.font.name = font
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    tf.anchor = anchor
    return box


def add_logo(slide, dark_bg=False):
    logo = LOGO_WHITE if dark_bg else LOGO_BLACK
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(0.1), Inches(5.2), height=Inches(0.3))


def main():
    # Open the extended template
    template_path = r"C:\Users\Christina\Code\Work\startuplab-partner-presentation\StartupLab_Extended_Template.pptx"
    prs = Presentation(template_path)
    blank = prs.slide_layouts[10]

    print(f"Starting with {len(prs.slides)} slides")

    # ========================================
    # SLIDE: 4-Up Stats (Summit Style)
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.5),
             "KEY FIGURES", size=15, font=FONT_HEADLINE, bold=True,
             color=COLORS['muted'], caps=True)

    add_text(slide, MARGIN, Inches(0.7), Inches(9.4), Inches(0.8),
             "STARTUPLAB SUMMIT 2025", size=32, font=FONT_HEADLINE,
             bold=True, color=COLORS['black'], caps=True)

    # 4 stat boxes
    stats = [
        ("1600+", "attendees"),
        ("140+", "corp leaders"),
        ("300+", "investors"),
        ("40+", "countries"),
    ]

    for i, (num, label) in enumerate(stats):
        x = MARGIN + (i * Inches(2.4))
        y = Inches(2.0)

        # Big number
        add_text(slide, x, y, Inches(2.2), Inches(1.2),
                 num, size=54, font=FONT_HEADLINE, bold=True,
                 color=COLORS['red'], align=PP_ALIGN.LEFT)

        # Label
        add_text(slide, x, y + Inches(1.1), Inches(2.2), Inches(0.5),
                 label, size=17, color=COLORS['black'])

    add_logo(slide)

    # ========================================
    # SLIDE: Portfolio Value Metrics
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.5),
             "PORTFOLIO IMPACT", size=15, font=FONT_HEADLINE, bold=True,
             color=COLORS['muted'], caps=True)

    add_text(slide, MARGIN, Inches(0.7), Inches(9.4), Inches(0.8),
             "VALUE CREATED SINCE 2012", size=32, font=FONT_HEADLINE,
             bold=True, color=COLORS['black'], caps=True)

    # Large metrics (3 columns)
    metrics = [
        ("8.1bn", "NOK", "combined revenue generated"),
        ("3,600+", "", "jobs created across portfolio"),
        ("18bn", "NOK", "top 10 companies valuation"),
    ]

    for i, (num, prefix, desc) in enumerate(metrics):
        x = MARGIN + (i * Inches(3.2))
        y = Inches(2.0)

        # Prefix (small, above number)
        if prefix:
            add_text(slide, x, y - Inches(0.3), Inches(2.8), Inches(0.3),
                     prefix, size=14, color=COLORS['muted'])

        # Big number
        add_text(slide, x, y, Inches(2.8), Inches(1.2),
                 num, size=48, font=FONT_HEADLINE, bold=True, color=COLORS['red'])

        # Description
        add_text(slide, x, y + Inches(1.0), Inches(2.8), Inches(0.8),
                 desc, size=15, color=COLORS['black'])

    add_logo(slide)

    # ========================================
    # SLIDE: Pipeline/Funnel Stats
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.8),
             "OUR PIPELINE", size=28, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    # Funnel visualization (4 rows)
    funnel = [
        ("1000+", "screened annually", Inches(9.0)),
        ("560", "startups since 2012", Inches(7.5)),
        ("<6%", "acceptance rate", Inches(6.0)),
        ("160+", "investments made", Inches(4.5)),
    ]

    for i, (num, label, width) in enumerate(funnel):
        y = Inches(1.3) + (i * Inches(1.0))
        x = (SLIDE_WIDTH - width) / 2  # Center

        # Bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.7))
        # Gradient from red to lighter
        opacity = 1.0 - (i * 0.15)
        r = int(0xFF * opacity + 0xF7 * (1 - opacity))
        g = int(0x33 * opacity + 0xF7 * (1 - opacity))
        b = int(0x33 * opacity + 0xF7 * (1 - opacity))
        set_fill(bar, RGBColor(r, g, b))
        bar.line.fill.background()

        # Number (left)
        add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(1.5), Inches(0.5),
                 num, size=24, font=FONT_HEADLINE, bold=True, color=COLORS['white'])

        # Label (right of number)
        add_text(slide, x + Inches(1.8), y + Inches(0.15), Inches(4), Inches(0.5),
                 label, size=16, color=COLORS['white'])

    add_logo(slide)

    # ========================================
    # SLIDE: Corporate-Startup Collab Canvas
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.2), Inches(9.4), Inches(0.5),
             "COLLABORATION FRAMEWORK", size=15, font=FONT_HEADLINE, bold=True,
             color=COLORS['muted'], caps=True)

    add_text(slide, MARGIN, Inches(0.5), Inches(9.4), Inches(0.5),
             "DEFINE CANVAS", size=24, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    # Canvas grid (2x2)
    canvas_items = [
        ("COLLABORATION OBJECTIVES", "What problem are we solving and why now?"),
        ("SUCCESS CRITERIA", "How will we measure success? KPIs and metrics"),
        ("STAKEHOLDERS", "Who needs to be involved? Champions and sponsors"),
        ("TIMELINE & RESOURCES", "What is the timeline? Budget and team allocation"),
    ]

    for i, (title, content) in enumerate(canvas_items):
        col = i % 2
        row = i // 2
        x = MARGIN + (col * Inches(4.8))
        y = Inches(1.2) + (row * Inches(2.0))

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(1.8))
        set_fill(box, COLORS['off_white'])
        box.line.color.rgb = COLORS['light_gray']

        # Title
        add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(4.2), Inches(0.4),
                 title, size=12, font=FONT_HEADLINE, bold=True, color=COLORS['red'], caps=True)

        # Content
        add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(4.2), Inches(1.2),
                 content, size=13, color=COLORS['muted'])

    add_logo(slide)

    # ========================================
    # SLIDE: Pilot Canvas
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.2), Inches(9.4), Inches(0.5),
             "COLLABORATION FRAMEWORK", size=15, font=FONT_HEADLINE, bold=True,
             color=COLORS['muted'], caps=True)

    add_text(slide, MARGIN, Inches(0.5), Inches(9.4), Inches(0.5),
             "PILOT CANVAS", size=24, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    canvas_items = [
        ("PILOT DESCRIPTION", "What exactly will be tested, and in what environment?"),
        ("PILOT METRICS", "Success criteria, KPIs, and data to collect"),
        ("RISKS & MITIGATIONS", "What could go wrong? Contingency plans"),
        ("NEXT STEPS", "Go/no-go criteria and scaling plan"),
    ]

    for i, (title, content) in enumerate(canvas_items):
        col = i % 2
        row = i // 2
        x = MARGIN + (col * Inches(4.8))
        y = Inches(1.2) + (row * Inches(2.0))

        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.5), Inches(1.8))
        set_fill(box, COLORS['off_white'])
        box.line.color.rgb = COLORS['light_gray']

        add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(4.2), Inches(0.4),
                 title, size=12, font=FONT_HEADLINE, bold=True, color=COLORS['red'], caps=True)

        add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(4.2), Inches(1.2),
                 content, size=13, color=COLORS['muted'])

    add_logo(slide)

    # ========================================
    # SLIDE: Case Study Card (Investment)
    # ========================================
    slide = prs.slides.add_slide(blank)

    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    set_fill(bg, COLORS['deep_black'])
    bg.line.fill.background()

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.4),
             "CASE STUDY", size=14, font=FONT_HEADLINE, color=COLORS['muted'], caps=True)

    add_text(slide, MARGIN, Inches(0.6), Inches(9.4), Inches(0.6),
             "CORPORATE × STARTUP", size=28, font=FONT_HEADLINE, bold=True,
             color=COLORS['white'], caps=True)

    # Company description
    add_text(slide, MARGIN, Inches(1.4), Inches(5.5), Inches(2.0),
             "Brief description of the collaboration, what problem was solved, "
             "and the outcomes achieved. Focus on measurable impact and key learnings.",
             size=16, color=COLORS['white'])

    # Metrics card (right side)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(1.4), Inches(3.4), Inches(3.2))
    set_fill(card, RGBColor(0x1A, 0x1A, 0x1A))
    card.line.fill.background()

    metrics = [
        ("Pilot Duration", "3 months"),
        ("Investment", "500K NOK"),
        ("ROI", "3.2x"),
        ("Status", "Scaling"),
    ]

    for i, (label, value) in enumerate(metrics):
        y = Inches(1.6) + (i * Inches(0.7))
        add_text(slide, Inches(6.5), y, Inches(3), Inches(0.3),
                 label, size=11, color=COLORS['muted'])
        add_text(slide, Inches(6.5), y + Inches(0.25), Inches(3), Inches(0.4),
                 value, size=18, font=FONT_HEADLINE, bold=True, color=COLORS['red'])

    add_logo(slide, dark_bg=True)

    # ========================================
    # SLIDE: Process Flow (Define → Scout → Align → Pilot)
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.6),
             "COLLABORATION JOURNEY", size=28, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    stages = [
        ("DEFINE", "Set objectives\n& success criteria"),
        ("SCOUT", "Find & evaluate\nstartup partners"),
        ("ALIGN", "Agree terms\n& expectations"),
        ("PILOT", "Test solution\n& measure results"),
    ]

    for i, (title, desc) in enumerate(stages):
        x = MARGIN + (i * Inches(2.4))
        y = Inches(1.8)

        # Circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(1.0), Inches(1.0))
        set_fill(circle, COLORS['red'])
        circle.line.fill.background()

        # Number in circle
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        # Arrow (except last)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + Inches(1.7), y + Inches(0.35),
                                           Inches(0.6), Inches(0.3))
            set_fill(arrow, COLORS['light_gray'])
            arrow.line.fill.background()

        # Title
        add_text(slide, x, y + Inches(1.2), Inches(2.2), Inches(0.4),
                 title, size=14, font=FONT_HEADLINE, bold=True,
                 color=COLORS['black'], align=PP_ALIGN.CENTER, caps=True)

        # Description
        add_text(slide, x, y + Inches(1.6), Inches(2.2), Inches(1.0),
                 desc, size=12, color=COLORS['muted'], align=PP_ALIGN.CENTER)

    add_logo(slide)

    # ========================================
    # SLIDE: Acceptance/Selectivity Stats
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.6),
             "SELECTIVITY & QUALITY", size=28, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    # Three stat cards
    cards = [
        ("<6%", "acceptance rate", "~680 applications in 2024"),
        ("560", "startups since 2012", "Norway's largest incubator"),
        ("100+", "tailored intros", "Corporate-startup matches"),
    ]

    for i, (num, title, sub) in enumerate(cards):
        x = MARGIN + (i * Inches(3.2))
        y = Inches(1.5)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.0), Inches(3.2))
        set_fill(card, COLORS['off_white'])
        card.line.fill.background()

        # Big number
        add_text(slide, x, y + Inches(0.5), Inches(3.0), Inches(1.2),
                 num, size=54, font=FONT_HEADLINE, bold=True,
                 color=COLORS['red'], align=PP_ALIGN.CENTER)

        # Title
        add_text(slide, x, y + Inches(1.6), Inches(3.0), Inches(0.5),
                 title, size=17, font=FONT_HEADLINE, bold=True,
                 color=COLORS['black'], align=PP_ALIGN.CENTER)

        # Subtitle
        add_text(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.6), Inches(0.8),
                 sub, size=13, color=COLORS['muted'], align=PP_ALIGN.CENTER)

    add_logo(slide)

    # ========================================
    # SLIDE: Long Text Slide with Bullets
    # ========================================
    slide = prs.slides.add_slide(blank)

    add_text(slide, MARGIN, Inches(0.3), Inches(9.4), Inches(0.6),
             "DETAILED CONTENT", size=28, font=FONT_HEADLINE, bold=True,
             color=COLORS['black'], caps=True)

    add_text(slide, MARGIN, Inches(0.9), Inches(9.4), Inches(0.4),
             "Subheading with additional context for the detailed information below",
             size=15, color=COLORS['muted'])

    # Two columns of bullet content
    left_content = """Key point one with detailed explanation that provides context and supporting information for the audience.

Key point two that builds on the previous point and adds new dimensions to the topic being discussed.

Key point three with specific examples or data that reinforce the main message."""

    right_content = """Additional context for point four that covers another aspect of the topic with relevant details.

Point five that addresses potential questions or concerns the audience might have.

Final point that ties everything together and leads to the next section or call to action."""

    add_text(slide, MARGIN, Inches(1.5), Inches(4.5), Inches(3.5),
             left_content, size=14, color=COLORS['black'])

    add_text(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(3.5),
             right_content, size=14, color=COLORS['black'])

    add_logo(slide)

    # Save
    prs.save(template_path)
    print(f"\nUpdated template saved with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
