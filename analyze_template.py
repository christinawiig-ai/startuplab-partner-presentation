"""
Analyze the existing StartupLab template to understand its structure
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def emu_to_inches(emu):
    """Convert EMU to inches"""
    return emu / 914400

def analyze_presentation(pptx_path):
    """Analyze a PowerPoint presentation structure"""
    prs = Presentation(pptx_path)

    print(f"\n{'='*60}")
    print(f"ANALYZING: {os.path.basename(pptx_path)}")
    print(f"{'='*60}")

    # Slide dimensions
    print(f"\nSlide Dimensions:")
    print(f"  Width: {emu_to_inches(prs.slide_width):.2f} inches")
    print(f"  Height: {emu_to_inches(prs.slide_height):.2f} inches")

    # Slide layouts
    print(f"\nSlide Layouts Available: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")

    # Analyze each slide
    print(f"\nTotal Slides: {len(prs.slides)}")

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'-'*50}")
        print(f"SLIDE {slide_num}")
        print(f"{'-'*50}")

        # Layout info
        if slide.slide_layout:
            print(f"Layout: {slide.slide_layout.name}")

        # Background
        if slide.background.fill.type is not None:
            fill = slide.background.fill
            if fill.type == 1:  # Solid
                try:
                    print(f"Background: Solid color")
                except:
                    pass

        # Shapes
        print(f"Shapes: {len(slide.shapes)}")

        for shape in slide.shapes:
            shape_info = f"  - {shape.shape_type}"

            # Position and size
            if hasattr(shape, 'left') and shape.left is not None:
                pos = f" @ ({emu_to_inches(shape.left):.1f}, {emu_to_inches(shape.top):.1f})"
                size = f" [{emu_to_inches(shape.width):.1f}x{emu_to_inches(shape.height):.1f}]"
                shape_info += pos + size

            # Text content
            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace('\n', ' ')
                if text.strip():
                    shape_info += f' "{text}..."' if len(shape.text_frame.text) > 50 else f' "{text}"'

                # Font info from first paragraph
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    if p.runs:
                        run = p.runs[0]
                        font_info = []
                        if run.font.name:
                            font_info.append(run.font.name)
                        if run.font.size:
                            font_info.append(f"{run.font.size.pt:.0f}pt")
                        if run.font.bold:
                            font_info.append("bold")
                        if font_info:
                            shape_info += f" ({', '.join(font_info)})"

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info += " [IMAGE]"

            print(shape_info)

# Analyze the official template
template_path = r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\Presentation Template (Updated_ 24.04.24).pptx"
if os.path.exists(template_path):
    analyze_presentation(template_path)
else:
    print(f"Template not found: {template_path}")

# Also analyze Value Pyramid template
template2_path = r"C:\Users\Christina\Code\Work\SL presentation engine\input\presentations\StartupLab - Value Pyramid TEMPLATE.pptx"
if os.path.exists(template2_path):
    analyze_presentation(template2_path)
